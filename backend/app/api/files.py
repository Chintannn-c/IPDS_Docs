from fastapi import APIRouter, UploadFile, File, Depends, HTTPException
from app.services.auth_service import oauth2_scheme, get_current_user
from app.services.encryption_service import encrypt_data, decrypt_data
from app.db.database import Database
from app.models.schemas import FileModel, LogActor, LogTarget
from app.services.scanner_engine_v2 import ScannerEngineV2 as ScannerEngine
from app.services.live_monitor import LiveMonitor
from app.websocket_manager import manager
from app.models.notification_models import NotificationCategory, NotificationPriority
from app.services.ai_service import AIService
import uuid
from datetime import datetime, timezone
import os
import fitz  # PyMuPDF
from PIL import Image
import pytesseract
from pptx import Presentation
import io
import hashlib

router = APIRouter()

# Configure Tesseract OCR path for Windows
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

UPLOAD_DIR = "uploads"
if not os.path.exists(UPLOAD_DIR):
    os.makedirs(UPLOAD_DIR)


# -------------------------------------------------------------
#  FILE UPLOAD (FULLY FIXED)
# -------------------------------------------------------------
@router.post("/upload", response_model=FileModel)
async def upload_file(
    file: UploadFile = File(...),
    current_user: dict = Depends(get_current_user)
):
    user_id = str(current_user["_id"])  # FIX: defined at top
    MAX_FILE_SIZE = 100 * 1024 * 1024  # 100MB
    content = await file.read()
    
    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(
            status_code=413,
            detail="File too large. Maximum size is 50MB."
        )

    # ------------------------------
    # DUPLICATE DETECTION
    # ------------------------------
    # Calculate SHA-256 hash of file content
    file_hash = hashlib.sha256(content).hexdigest()
    
    # Check if file with same hash already exists for this user
    db = Database.get_db()
    existing_file = db.files.find_one({
        "owner_id": current_user["_id"],
        "file_hash": file_hash
    })
    
    if existing_file:
        raise HTTPException(
            status_code=409,
            detail=f"Duplicate file detected. '{existing_file['filename']}' already exists in your storage."
        )

    # ------------------------------
    # 1. SCAN FILE USING IPDS ENGINE
    # ------------------------------
    scan_result = ScannerEngine.scan_file(content, filename=file.filename, mime_type=file.content_type)

    safety_score = scan_result["safety_score"]
    findings = scan_result["findings"]
    recommendation = scan_result.get("recommendation", "No issues")

    # ------------------------------
    # 2. IPDS PREVENTION LAYER (MODIFIED)
    # ------------------------------
    is_risky = safety_score < 70
    is_quarantined = safety_score < 50
    risk_reason = recommendation if is_risky else None

    # We now allow the upload but flag it
    if is_quarantined:
        print(f"⚠️ QUARANTINE: File {file.filename} (Score: {safety_score}) flagged for security review.")

    # -------------------------------------------------
    # 3. SAFE FILE → ENCRYPT, SAVE PHYSICALLY & IN DB
    # -------------------------------------------------
    encrypted = encrypt_data(content)

    file_id = str(uuid.uuid4())
    file_path = os.path.join(UPLOAD_DIR, file_id)

    with open(file_path, "wb") as f:
        f.write(encrypted)

    db = Database.get_db()

    new_file = {
        "_id": file_id,
        "filename": file.filename,
        "size": len(content),
        "upload_date": datetime.utcnow(),
        "owner_id": current_user["_id"],
        "path": file_path,
        "file_hash": file_hash,  # Store hash for duplicate detection
        "safety_score": safety_score,
        "scan_status": "scanned",
        "scan_findings": findings,
        "is_risky": is_risky,
        "is_quarantined": is_quarantined,
        "risk_reason": risk_reason
    }

    db.files.insert_one(new_file)

    # -----------------------------------------
    # 4. LiveMonitor — track file & log activity
    # -----------------------------------------
    await LiveMonitor.init_file_tracking(file_id, user_id)

    actor = LogActor(
        user_id=user_id,
        name=current_user["name"],
        role=current_user.get("role", "USER"),
    )
    target = LogTarget(type="FILE", id=file_id, name=file.filename)

    # Determine log status based on safety score
    # Score < 50: Dangerous (but allowed through), Score < 70: Risky/Caution
    if safety_score < 50:
        log_status = "DANGER"
        log_action = "Risky File Uploaded"
    elif safety_score < 70:
        log_status = "WARNING"
        log_action = "File Uploaded (Caution)"
    else:
        log_status = "SUCCESS"
        log_action = "File Uploaded"

    await LiveMonitor.log_activity(
        actor=actor,
        action=log_action,
        status=log_status,
        target=target,
        metadata={
            "size": new_file["size"],
            "safety_score": safety_score,
            "findings_count": len(findings),
        },
    )
    
    # Send notification about successful upload
    notification_title = " File Uploaded" if safety_score >= 70 else "File Uploaded (Low Safety Score)"
    notification_priority = NotificationPriority.MEDIUM.value if safety_score >= 70 else NotificationPriority.HIGH.value
    
    await manager.create_and_broadcast_notification(
        user_id=user_id,
        title=notification_title,
        message=f"'{file.filename}' uploaded successfully. Safety score: {safety_score}%",
        category=NotificationCategory.FILE.value,
        priority=notification_priority,
        data={
            "filename": file.filename,
            "file_id": file_id,
            "safety_score": safety_score,
            "size": new_file["size"]
        }
    )

    # -----------------------------------------
    # 5. API RESPONSE
    # -----------------------------------------
    return FileModel(
        id=file_id,
        filename=new_file["filename"],
        size=new_file["size"],
        upload_date=new_file["upload_date"],
        owner_id=new_file["owner_id"],
        safety_score=safety_score,
        scan_status="scanned",
        is_risky=is_risky,
        is_quarantined=is_quarantined,
        risk_reason=risk_reason
    )


# -------------------------------------------------------------
#  FILE DOWNLOAD
# -------------------------------------------------------------
@router.get("/download/{file_id}")
async def download_file(file_id: str, current_user: dict = Depends(get_current_user)):
    db = Database.get_db()
    file_record = db.files.find_one({"_id": file_id, "owner_id": current_user["_id"]})

    if not file_record:
        raise HTTPException(status_code=404, detail="File not found")

    with open(file_record["path"], "rb") as f:
        encrypted = f.read()

    decrypted = decrypt_data(encrypted)

    from fastapi.responses import Response
    return Response(
        content=decrypted,
        media_type="application/octet-stream",
        headers={"Content-Disposition": f"attachment; filename={file_record['filename']}"}
    )


# -------------------------------------------------------------
#  LIST FILES
# -------------------------------------------------------------
@router.get("/", response_model=list[FileModel])
async def list_files(current_user: dict = Depends(get_current_user)):
    db = Database.get_db()
    # Sort by upload_date descending (newest first)
    files = list(db.files.find({"owner_id": current_user["_id"]}).sort("upload_date", -1))

    return [
        FileModel(
            id=str(f["_id"]),
            filename=f["filename"],
            size=f["size"],
            upload_date=f["upload_date"],
            owner_id=f["owner_id"],
            safety_score=f.get("safety_score", 100),
            scan_status=f.get("scan_status", "scanned"),
            is_risky=f.get("is_risky", False),
            is_quarantined=f.get("is_quarantined", False),
            risk_reason=f.get("risk_reason")
        )
        for f in files
    ]


# -------------------------------------------------------------
#  DELETE FILE
# -------------------------------------------------------------
@router.delete("/{file_id}")
async def delete_file(file_id: str, current_user: dict = Depends(get_current_user)):
    db = Database.get_db()
    file_record = db.files.find_one({"_id": file_id, "owner_id": current_user["_id"]})

    if not file_record:
        raise HTTPException(status_code=404, detail="File not found")

    if os.path.exists(file_record["path"]):
        os.remove(file_record["path"])

    db.files.delete_one({"_id": file_id})
    
    # Log the action (userId, fileId, timestamp, risk reason)
    actor = LogActor(
        user_id=str(current_user["_id"]),
        name=current_user["name"],
        role=current_user.get("role", "USER"),
    )
    target = LogTarget(type="FILE", id=file_id, name=file_record["filename"])
    
    await LiveMonitor.log_activity(
        actor=actor,
        action="File Deleted",
        status="SUCCESS",
        target=target,
        metadata={
            "filename": file_record["filename"],
            "was_risky": file_record.get("is_risky", False),
            "risk_reason": file_record.get("risk_reason")
        }
    )

    # Send notification about file deletion
    await manager.create_and_broadcast_notification(
        user_id=str(current_user["_id"]),
        title="🗑️ File Deleted",
        message=f"'{file_record['filename']}' has been deleted.",
        category=NotificationCategory.FILE.value,
        priority=NotificationPriority.LOW.value,
        data={"filename": file_record["filename"], "file_id": file_id}
    )

    return {"message": "File deleted successfully"}

# -------------------------------------------------------------
#  CONFIRM/DECLINE RISK
# -------------------------------------------------------------
@router.post("/{file_id}/confirm-risk")
async def confirm_file_risk(file_id: str, current_user: dict = Depends(get_current_user)):
    """User confirms removal of a risky file."""
    db = Database.get_db()
    file_record = db.files.find_one({"_id": file_id, "owner_id": current_user["_id"]})

    if not file_record:
        raise HTTPException(status_code=404, detail="File not found")

    # Perform full deletion
    if os.path.exists(file_record["path"]):
        os.remove(file_record["path"])

    db.files.delete_one({"_id": file_id})

    # Log the action
    actor = LogActor(
        user_id=str(current_user["_id"]),
        name=current_user["name"],
        role=current_user.get("role", "USER"),
    )
    target = LogTarget(type="FILE", id=file_id, name=file_record["filename"])
    
    await LiveMonitor.log_activity(
        actor=actor,
        action="Risky File Removal Confirmed",
        status="SUCCESS",
        target=target,
        metadata={
            "filename": file_record["filename"],
            "risk_reason": file_record.get("risk_reason")
        }
    )

    return {"message": "Risky file removed as per user confirmation"}

@router.post("/{file_id}/decline-risk")
async def decline_file_risk(file_id: str, current_user: dict = Depends(get_current_user)):
    """User accepts the risk - file keeps its risk status but is no longer quarantined."""
    db = Database.get_db()
    file_record = db.files.find_one({"_id": file_id, "owner_id": current_user["_id"]})

    if not file_record:
        raise HTTPException(status_code=404, detail="File not found")

    # Remove quarantine but KEEP original risk status and safety score
    db.files.update_one(
        {"_id": file_id},
        {
            "$set": {
                "is_quarantined": False,  # Remove quarantine
                "risk_accepted": True,     # Mark as accepted
                "risk_accepted_at": datetime.now(timezone.utc).isoformat(),
                "risk_accepted_by": str(current_user["_id"])
                # NOTE: is_risky, safety_score, and risk_reason remain UNCHANGED
            }
        }
    )

    # Log the risk acceptance
    actor = LogActor(
        user_id=str(current_user["_id"]),
        name=current_user["name"],
        role=current_user.get("role", "USER"),
    )
    target = LogTarget(type="FILE", id=file_id, name=file_record["filename"])
    
    await LiveMonitor.log_activity(
        actor=actor,
        action=f"Risk Accepted - {file_record['filename']} (Score: {file_record.get('safety_score', 'N/A')})",
        status="WARNING",  # Show as warning since file is still risky
        target=target,
        metadata={
            "filename": file_record["filename"],
            "risk_reason": file_record.get("risk_reason"),
            "safety_score": file_record.get("safety_score"),
            "is_risky": True,  # File remains risky
            "risk_accepted": True
        }
    )

    # Send notification
    await manager.create_and_broadcast_notification(
        user_id=str(current_user["_id"]),
        title="Risk Accepted",
        message=f"Risk accepted for '{file_record['filename']}'. File remains risky but is accessible.",
        category=NotificationCategory.SECURITY.value,
        priority=NotificationPriority.MEDIUM.value,
        data={"filename": file_record["filename"], "file_id": file_id}
    )

    return {"message": "Risk accepted. File remains risky but is now accessible."}

# -------------------------------------------------------------
#  DOCUMENT ANALYSIS (NEW)
# -------------------------------------------------------------
def extract_text_from_file(file_path: str, filename: str) -> dict:
    """
    Extracts text from various file formats including PDF, images, and PowerPoint.
    Returns a dict with text and metadata about extraction method.
    """
    ext = os.path.splitext(filename)[1].lower()
    
    # Decrypt first
    with open(file_path, "rb") as f:
        encrypted_data = f.read()
    
    decrypted_data = decrypt_data(encrypted_data)
    
    ocr_executed = False
    text_source = "none"
    extracted_text = ""
    
    # PDF extraction with OCR fallback
    if ext == ".pdf":
        try:
            # Step 1: Try standard PDF text extraction
            doc = fitz.open(stream=decrypted_data, filetype="pdf")
            text = ""
            for page in doc:
                text += page.get_text()
            doc.close()
            
            # Step 2: Check if extraction returned meaningful text
            if text.strip():  # Has readable text
                extracted_text = text
                text_source = "parser"
            else:
                # Step 3: Text extraction failed - trigger OCR
                print(f"DEBUG: PDF text extraction returned empty. Triggering OCR for {filename}...")
                ocr_executed = True
                
                # Parallel OCR processing (thread-safe)
                def process_page_ocr(page_num, pdf_bytes):
                    """Each thread gets its own document instance (thread-safe)."""
                    try:
                        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
                        page = doc[page_num]
                        
                        # Render page to image (higher DPI for better OCR)
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom
                        img_data = pix.tobytes("png")
                        
                        # Run OCR on the image
                        image = Image.open(io.BytesIO(img_data))
                        page_text = pytesseract.image_to_string(image)
                        
                        doc.close()  # Clean up this thread's document
                        return page_num, page_text
                    except Exception as e:
                        print(f"[OCR] Error processing page {page_num}: {e}")
                        return page_num, ""
                
                # Get page count safely
                temp_doc = fitz.open(stream=decrypted_data, filetype="pdf")
                page_count = len(temp_doc)
                temp_doc.close()
                
                # Process pages in parallel (up to 4 concurrent)
                from concurrent.futures import ThreadPoolExecutor, as_completed
                max_workers = min(4, page_count)
                
                ocr_results = {}
                with ThreadPoolExecutor(max_workers=max_workers) as executor:
                    futures = {executor.submit(process_page_ocr, i, decrypted_data): i for i in range(page_count)}
                    for future in as_completed(futures):
                        page_num, page_text = future.result()
                        ocr_results[page_num] = page_text
                        print(f"[OCR] Completed page {page_num + 1}/{page_count}")
                
                # Reconstruct text in correct order
                ocr_text = ""
                for i in range(page_count):
                    ocr_text += f"\n--- Page {i + 1} ---\n{ocr_results.get(i, '')}"
                
                if ocr_text.strip():
                    extracted_text = ocr_text
                    text_source = "ocr"
                    print(f"DEBUG: Parallel OCR successful. Extracted {len(ocr_text)} characters.")
                else:
                    print(f"DEBUG: OCR also returned empty. Document may be image-only or encrypted.")
                    text_source = "none"
                    
        except Exception as e:
            print(f"Error extracting PDF text: {e}")
            text_source = "none"
    
    # Text file extraction
    elif ext in [".txt", ".md", ".py", ".dart", ".js", ".html", ".css"]:
        try:
            extracted_text = decrypted_data.decode("utf-8")
            text_source = "parser"
        except UnicodeDecodeError:
            try:
                extracted_text = decrypted_data.decode("latin-1")
                text_source = "parser"
            except:
                text_source = "none"
    
    # Image extraction (OCR)
    elif ext in [".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".webp"]:
        try:
            image = Image.open(io.BytesIO(decrypted_data))
            extracted_text = pytesseract.image_to_string(image)
            text_source = "ocr"
            ocr_executed = True
        except Exception as e:
            print(f"Error extracting text from image: {e}")
            text_source = "none"
    
    # PowerPoint extraction
    elif ext in [".pptx", ".ppt"]:
        try:
            prs = Presentation(io.BytesIO(decrypted_data))
            text = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"\n--- Slide {slide_num} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            extracted_text = text
            text_source = "parser"
        except Exception as e:
            print(f"Error extracting PowerPoint text: {e}")
            text_source = "none"
    
    return {
        "text": extracted_text.strip(),
        "text_source": text_source,
        "ocr_executed": ocr_executed
    }

@router.post("/{file_id}/analyze")
async def analyze_file(file_id: str, current_user: dict = Depends(get_current_user)):
    db = Database.get_db()
    file_record = db.files.find_one({"_id": file_id, "owner_id": current_user["_id"]})

    if not file_record:
        print(f"DEBUG: File {file_id} not found for user {current_user['_id']}")
        raise HTTPException(status_code=404, detail="File not found")

    # Extract text with metadata
    print(f"DEBUG: Extracting text from {file_record['filename']}...")
    extraction_result = extract_text_from_file(file_record["path"], file_record["filename"])
    text = extraction_result["text"]
    text_source = extraction_result["text_source"]
    ocr_executed = extraction_result["ocr_executed"]
    print(f"DEBUG: Extraction complete. Source: {text_source}, OCR: {ocr_executed}, Length: {len(text)}")
    
    if not text or not text.strip():
        print(f"DEBUG: Extraction failed or returned empty text for {file_record['filename']}")
        raise HTTPException(status_code=400, detail="Document contains no extractable text or format unsupported.")

    print(f"DEBUG: Successfully extracted {len(text)} characters.")

    # Process with chunking support
    analysis = AIService.analyze_document_chunked(text)
    
    # Add extraction metadata
    analysis["text_source"] = text_source
    analysis["ocr_executed"] = ocr_executed
    
    # Store initial assessment in document_summaries as v1
    # We do this to unify the versioning system
    summary_doc = {
        "document_id": file_id,
        "document_name": file_record["filename"],
        "summary": str(analysis.get("summary", "") or ""),
        "key_points": analysis.get("key_points", []) or [],
        "risk_flags": analysis.get("risk_flags", []) or [],
        "content_preview": str(analysis.get("content_preview", "") or ""),
        "version": 1,
        "summarized_by": {
            "user_id": str(current_user["_id"]),
            "user_email": current_user["email"]
        },
        "summarized_at": datetime.now(timezone.utc),
        "ai_model": "mistral-large-latest"
    }
    
    print(f"[FILES] Created initial summary v1 for: {file_record['filename']}")
    db.document_summaries.insert_one(summary_doc)

    # Store analysis in DB (main files record for fast check/legacy)
    db.files.update_one(
        {"_id": file_id},
        {"$set": {"analysis": analysis, "analyzed_at": datetime.now(timezone.utc)}}
    )

    # Log activity
    actor = LogActor(
        user_id=str(current_user["_id"]),
        name=current_user["name"],
        role=current_user.get("role", "USER"),
    )
    target = LogTarget(type="FILE", id=file_id, name=file_record["filename"])
    
    await LiveMonitor.log_activity(
        actor=actor,
        action="Document Analyzed",
        status="SUCCESS",
        target=target,
        metadata={"file_id": file_id}
    )

    return analysis

@router.get("/{file_id}/analysis")
async def get_file_analysis(file_id: str, current_user: dict = Depends(get_current_user)):
    db = Database.get_db()
    
    # Check if a versioned summary exists (this is the new standard)
    latest_summary = db.document_summaries.find_one(
        {"document_id": file_id},
        sort=[("version", -1)]
    )
    
    if latest_summary:
        # Map it to the format expected by the frontend UI
        return {
            "summary": latest_summary.get("summary", ""),
            "key_points": latest_summary.get("key_points", []),
            "risk_flags": latest_summary.get("risk_flags", []),
            "content_preview": latest_summary.get("content_preview", ""),
            "version": latest_summary.get("version", 1),
            "_id": str(latest_summary["_id"])
        }

    # Fallback to legacy file analysis if no versioned history exists
    file_record = db.files.find_one({"_id": file_id, "owner_id": current_user["_id"]})
    if not file_record:
        raise HTTPException(status_code=404, detail="File not found")
    
    if "analysis" not in file_record:
        raise HTTPException(status_code=404, detail="No analysis found for this file. Please analyze first.")

    return file_record["analysis"]
